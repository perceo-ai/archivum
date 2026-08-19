"""File/URL parsers. Each returns a ParsedDoc dataclass."""

from __future__ import annotations

import asyncio
import base64
import csv
import email
import io
import json
import logging
import mailbox
import os
import re
import subprocess
import tempfile
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any
from urllib.parse import unquote, urlparse

import httpx

logger = logging.getLogger(__name__)

# ── Lazy-loaded optional dependencies ────────────────────────────────────────

_whisper_model = None  # cached whisper model (loading is slow)

try:
    import whisper as _whisper_lib  # type: ignore[import]
    _whisper_available = True
except ImportError:
    _whisper_lib = None
    _whisper_available = False


def _get_whisper_model():
    global _whisper_model
    if not _whisper_available:
        raise UnsupportedFileTypeError(
            "Audio/video transcription is not enabled yet. Open Settings and run Install / Enable under Audio Transcription."
        )
    if _whisper_model is None:
        _whisper_model = _whisper_lib.load_model("base")
    return _whisper_model


# ── Data model ────────────────────────────────────────────────────────────────

@dataclass
class ParsedDoc:
    text: str
    metadata: dict[str, Any] = field(default_factory=dict)
    source: str = ""


class UnsupportedFileTypeError(ValueError):
    """Raised when the file extension is not handled."""


# ── Helpers ───────────────────────────────────────────────────────────────────

def _strip_rst_directives(text: str) -> str:
    """Remove RST directives (.. foo:: ) and field lists."""
    text = re.sub(r"\.\. \w[\w-]*::.*", "", text)
    text = re.sub(r"^\s*:\w[\w -]*:.*$", "", text, flags=re.MULTILINE)
    return text


def _html_to_text(html_content: str, url: str = "") -> str:
    """Extract readable text from HTML, preferring main content."""
    try:
        from readability import Document
        doc = Document(html_content)
        main_html = doc.summary()
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(main_html, "html.parser")
        return soup.get_text(separator="\n", strip=True)
    except Exception:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(html_content, "html.parser")
        # Remove script/style/nav tags
        for tag in soup(["script", "style", "nav", "header", "footer", "aside"]):
            tag.decompose()
        return soup.get_text(separator="\n", strip=True)


def _message_text(msg: email.message.Message) -> str:
    if msg.is_multipart():
        for part in msg.walk():
            if part.get_content_type() == "text/plain":
                payload = part.get_payload(decode=True)
                if payload:
                    return payload.decode("utf-8", errors="replace")
        return ""

    payload = msg.get_payload(decode=True)
    if payload:
        return payload.decode("utf-8", errors="replace")
    payload_text = msg.get_payload()
    return payload_text if isinstance(payload_text, str) else ""


def _format_email_message(msg: email.message.Message) -> str:
    return "\n".join(
        [
            f"From: {msg.get('From', '')}",
            f"To: {msg.get('To', '')}",
            f"Subject: {msg.get('Subject', '')}",
            f"Date: {msg.get('Date', '')}",
            "",
            _message_text(msg),
        ]
    ).strip()


def _compact_blank_lines(text: str) -> str:
    return re.sub(r"\n{3,}", "\n\n", text).strip()


def _rtf_to_text(raw: str) -> str:
    """Best-effort RTF text extraction without adding a heavyweight dependency."""
    text = raw.replace("\\par", "\n").replace("\\line", "\n")
    text = re.sub(r"\\'[0-9a-fA-F]{2}", " ", text)
    text = re.sub(r"\\[a-zA-Z]+-?\d* ?", "", text)
    text = text.replace("{", "").replace("}", "")
    text = text.replace("\\", "")
    return _compact_blank_lines(re.sub(r"[ \t]+", " ", text))


def _xml_to_text(raw: str) -> str:
    try:
        from bs4 import BeautifulSoup

        soup = BeautifulSoup(raw, "xml")
        return soup.get_text(separator="\n", strip=True)
    except Exception:
        return _html_to_text(raw)


def _format_timestamp(seconds: Any) -> str:
    try:
        total = max(0, int(float(seconds)))
    except (TypeError, ValueError):
        total = 0
    hours, remainder = divmod(total, 3600)
    minutes, secs = divmod(remainder, 60)
    return f"{hours:02d}:{minutes:02d}:{secs:02d}"


def _format_transcription(result: dict[str, Any]) -> str:
    segments = result.get("segments")
    if isinstance(segments, list) and segments:
        lines = []
        for segment in segments:
            if not isinstance(segment, dict):
                continue
            text = str(segment.get("text") or "").strip()
            if not text:
                continue
            lines.append(
                "["
                f"{_format_timestamp(segment.get('start'))} - "
                f"{_format_timestamp(segment.get('end'))}"
                f"] {text}"
            )
        if lines:
            return "\n".join(lines)
    return str(result.get("text") or "").strip()


def _transcription_metadata(result: dict[str, Any], media_type: str) -> dict[str, Any]:
    metadata: dict[str, Any] = {"type": media_type}
    if result.get("language"):
        metadata["language"] = result["language"]
    if result.get("duration") is not None:
        metadata["duration_seconds"] = result.get("duration")
    segments = result.get("segments")
    if isinstance(segments, list):
        metadata["segments"] = len(segments)
    return metadata


def _parse_zip_archive(path: Path) -> ParsedDoc:
    parsed_parts: list[str] = []
    skipped: list[str] = []

    with tempfile.TemporaryDirectory(prefix="archivum-archive-") as tmp:
        tmp_dir = Path(tmp)
        with zipfile.ZipFile(path) as archive:
            for info in archive.infolist():
                member_name = info.filename
                if (
                    info.is_dir()
                    or member_name.startswith("__MACOSX/")
                    or Path(member_name).name.startswith(".")
                ):
                    continue
                safe_name = Path(member_name).name
                if not safe_name:
                    continue
                member_path = tmp_dir / safe_name
                member_path.write_bytes(archive.read(info))
                try:
                    doc = parse_file(member_path)
                except UnsupportedFileTypeError:
                    skipped.append(member_name)
                    continue
                title = member_name
                parsed_parts.append(f"## {title}\n\n{doc.text.strip()}")

    if not parsed_parts and skipped:
        raise UnsupportedFileTypeError(
            f"No supported files found in archive {path.name}; skipped {len(skipped)} members"
        )

    text = "\n\n---\n\n".join(part for part in parsed_parts if part.strip())
    if skipped:
        text = f"{text}\n\nSkipped unsupported archive members: {', '.join(skipped[:20])}".strip()

    return ParsedDoc(
        text=text,
        source=str(path),
        metadata={
            "type": "archive",
            "format": "zip",
            "files_parsed": len(parsed_parts),
            "files_skipped": len(skipped),
        },
    )


def _download_filename(url: str, content_type: str, content_disposition: str) -> str:
    match = re.search(
        r'filename\*?=(?:UTF-8\'\')?"?([^";]+)"?',
        content_disposition,
        flags=re.I,
    )
    if match:
        return Path(unquote(match.group(1))).name

    url_name = Path(unquote(urlparse(url).path)).name
    if url_name and Path(url_name).suffix:
        return url_name

    content_type = content_type.split(";", 1)[0].strip().lower()
    suffix_by_type = {
        "application/pdf": ".pdf",
        "application/epub+zip": ".epub",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": ".xlsx",
        "application/zip": ".zip",
        "text/rtf": ".rtf",
        "application/rtf": ".rtf",
        "application/xml": ".xml",
        "text/xml": ".xml",
    }
    return f"download{suffix_by_type.get(content_type, '.bin')}"


# ── Extension dispatch ────────────────────────────────────────────────────────

def parse_file(path: Path) -> ParsedDoc:
    """Synchronous file parser. Call from executor for async contexts."""
    suffix = path.suffix.lower()

    # ── Plain text / Markdown / RST ──────────────────────────────────────────
    if suffix in {".md", ".txt", ".rst", ".text", ".log"}:
        text = path.read_text(encoding="utf-8", errors="replace")
        if suffix == ".rst":
            text = _strip_rst_directives(text)
        return ParsedDoc(text=text, source=str(path), metadata={"type": suffix.lstrip(".")})

    # ── RTF / XML ────────────────────────────────────────────────────────────
    if suffix == ".rtf":
        return ParsedDoc(
            text=_rtf_to_text(path.read_text(encoding="utf-8", errors="replace")),
            source=str(path),
            metadata={"type": "rtf"},
        )

    if suffix in {".xml", ".rss", ".atom"}:
        return ParsedDoc(
            text=_xml_to_text(path.read_text(encoding="utf-8", errors="replace")),
            source=str(path),
            metadata={"type": "xml"},
        )

    # ── PDF ──────────────────────────────────────────────────────────────────
    if suffix == ".pdf":
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(str(path))
            pages_text = []
            for page_num, page in enumerate(doc):
                pages_text.append(f"[Page {page_num + 1}]\n{page.get_text()}")
            doc.close()
            return ParsedDoc(
                text="\n\n".join(pages_text),
                source=str(path),
                metadata={"type": "pdf", "pages": len(pages_text)},
            )
        except ImportError:
            raise UnsupportedFileTypeError("PyMuPDF not installed — cannot parse PDF")

    # ── HTML ─────────────────────────────────────────────────────────────────
    if suffix in {".html", ".htm"}:
        html = path.read_text(encoding="utf-8", errors="replace")
        return ParsedDoc(
            text=_html_to_text(html),
            source=str(path),
            metadata={"type": "html"},
        )

    # ── DOCX ─────────────────────────────────────────────────────────────────
    if suffix == ".docx":
        try:
            from docx import Document
            doc = Document(str(path))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return ParsedDoc(
                text="\n\n".join(paragraphs),
                source=str(path),
                metadata={"type": "docx"},
            )
        except ImportError:
            raise UnsupportedFileTypeError("python-docx not installed")

    # ── PPTX ─────────────────────────────────────────────────────────────────
    if suffix == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(str(path))
            slides_text = []
            for i, slide in enumerate(prs.slides):
                slide_parts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                slide_parts.append(t)
                if slide_parts:
                    slides_text.append(f"[Slide {i + 1}]\n" + "\n".join(slide_parts))
            return ParsedDoc(
                text="\n\n".join(slides_text),
                source=str(path),
                metadata={"type": "pptx", "slides": len(slides_text)},
            )
        except ImportError:
            raise UnsupportedFileTypeError("python-pptx not installed")

    # ── Excel / CSV ───────────────────────────────────────────────────────────
    if suffix in {".xlsx", ".xls"}:
        try:
            import pandas as pd
            xl = pd.ExcelFile(str(path))
            parts = []
            for sheet_name in xl.sheet_names:
                df = xl.parse(sheet_name)
                parts.append(
                    f"## Sheet: {sheet_name}\n"
                    f"Shape: {df.shape[0]} rows × {df.shape[1]} columns\n"
                    f"Columns: {', '.join(str(c) for c in df.columns)}\n\n"
                    f"First 10 rows:\n{df.head(10).to_markdown(index=False)}\n\n"
                    f"Summary stats:\n{df.describe(include='all').to_markdown()}"
                )
            return ParsedDoc(
                text="\n\n".join(parts),
                source=str(path),
                metadata={"type": "xlsx"},
            )
        except ImportError:
            raise UnsupportedFileTypeError("pandas/openpyxl not installed")

    if suffix == ".csv":
        try:
            import pandas as pd
            df = pd.read_csv(str(path))
            text = (
                f"CSV: {path.name}\n"
                f"Shape: {df.shape[0]} rows × {df.shape[1]} columns\n"
                f"Columns: {', '.join(str(c) for c in df.columns)}\n\n"
                f"First 10 rows:\n{df.head(10).to_markdown(index=False)}\n\n"
                f"Summary stats:\n{df.describe(include='all').to_markdown()}"
            )
            return ParsedDoc(text=text, source=str(path), metadata={"type": "csv"})
        except ImportError:
            # Fallback to stdlib csv
            with open(path, newline="", encoding="utf-8", errors="replace") as f:
                reader = csv.reader(f)
                rows = list(reader)
            header = rows[0] if rows else []
            text = f"CSV: {path.name}\nColumns: {', '.join(header)}\nRows: {len(rows) - 1}\n\n"
            text += "\n".join(",".join(r) for r in rows[:11])
            return ParsedDoc(text=text, source=str(path), metadata={"type": "csv"})

    # ── JSON / JSONL ──────────────────────────────────────────────────────────
    if suffix == ".json":
        raw = path.read_text(encoding="utf-8", errors="replace")
        try:
            obj = json.loads(raw)
            if isinstance(obj, list):
                sample = obj[:3]
                summary = f"JSON array with {len(obj)} items.\nSample:\n{json.dumps(sample, indent=2)}"
            elif isinstance(obj, dict):
                keys = list(obj.keys())
                summary = f"JSON object with keys: {', '.join(str(k) for k in keys[:20])}\n\nContent:\n{json.dumps(obj, indent=2)[:4000]}"
            else:
                summary = raw[:4000]
        except json.JSONDecodeError:
            summary = raw[:4000]
        return ParsedDoc(text=summary, source=str(path), metadata={"type": "json"})

    if suffix == ".jsonl":
        lines = path.read_text(encoding="utf-8", errors="replace").splitlines()
        sample = []
        for line in lines[:5]:
            try:
                sample.append(json.loads(line))
            except Exception:
                pass
        text = f"JSONL: {len(lines)} records\nSample (first 5):\n{json.dumps(sample, indent=2)}"
        return ParsedDoc(text=text, source=str(path), metadata={"type": "jsonl"})

    # ── Archives ─────────────────────────────────────────────────────────────
    if suffix == ".zip":
        return _parse_zip_archive(path)

    # ── EPUB ─────────────────────────────────────────────────────────────────
    if suffix == ".epub":
        try:
            import ebooklib
            from ebooklib import epub
            from bs4 import BeautifulSoup

            book = epub.read_epub(str(path))
            chapters = []
            for item in book.get_items():
                if item.get_type() == ebooklib.ITEM_DOCUMENT:
                    soup = BeautifulSoup(item.get_content(), "html.parser")
                    text = soup.get_text(separator="\n", strip=True)
                    if text.strip():
                        chapters.append(text)
            return ParsedDoc(
                text="\n\n---\n\n".join(chapters),
                source=str(path),
                metadata={"type": "epub", "chapters": len(chapters)},
            )
        except ImportError:
            raise UnsupportedFileTypeError("ebooklib not installed")

    # ── Code files ────────────────────────────────────────────────────────────
    code_extensions = {
        ".py": "python",
        ".js": "javascript",
        ".ts": "typescript",
        ".go": "go",
        ".rs": "rust",
        ".sh": "bash",
        ".bash": "bash",
        ".zsh": "zsh",
        ".rb": "ruby",
        ".java": "java",
        ".c": "c",
        ".cpp": "cpp",
        ".h": "c",
        ".hpp": "cpp",
        ".kt": "kotlin",
        ".swift": "swift",
        ".php": "php",
        ".sql": "sql",
        ".yaml": "yaml",
        ".yml": "yaml",
        ".toml": "toml",
        ".ini": "ini",
        ".cfg": "ini",
        ".jsonc": "jsonc",
        ".css": "css",
        ".scss": "scss",
        ".sass": "sass",
        ".less": "less",
        ".jsx": "jsx",
        ".tsx": "tsx",
        ".mjs": "javascript",
        ".cjs": "javascript",
        ".vue": "vue",
        ".svelte": "svelte",
        ".mdx": "mdx",
        ".dockerfile": "dockerfile",
    }
    if suffix in code_extensions:
        lang = code_extensions[suffix]
        code = path.read_text(encoding="utf-8", errors="replace")
        return ParsedDoc(
            text=f"```{lang}\n{code}\n```",
            source=str(path),
            metadata={"type": "code", "language": lang},
        )

    # ── SRT / VTT subtitle files ──────────────────────────────────────────────
    if suffix in {".srt", ".vtt"}:
        raw = path.read_text(encoding="utf-8", errors="replace")
        # Strip timestamp lines and sequence numbers
        lines = []
        for line in raw.splitlines():
            line = line.strip()
            # Skip sequence numbers, timestamps, and WEBVTT header
            if re.match(r"^\d+$", line):
                continue
            if re.match(r"[\d:,\. ]+-->", line):
                continue
            if line in {"WEBVTT", ""}:
                continue
            lines.append(line)
        return ParsedDoc(
            text=" ".join(lines),
            source=str(path),
            metadata={"type": suffix.lstrip(".")},
        )

    # ── EML ───────────────────────────────────────────────────────────────────
    if suffix == ".eml":
        raw = path.read_bytes()
        msg = email.message_from_bytes(raw)
        return ParsedDoc(
            text=_format_email_message(msg),
            source=str(path),
            metadata={"type": "eml", "subject": msg.get("Subject", "")},
        )

    if suffix == ".mbox":
        messages = mailbox.mbox(str(path), create=False)
        parts = [_format_email_message(msg) for msg in messages]
        return ParsedDoc(
            text="\n\n---\n\n".join(part for part in parts if part),
            source=str(path),
            metadata={"type": "mbox", "messages": len(parts)},
        )

    # ── Images ────────────────────────────────────────────────────────────────
    if suffix in {".png", ".jpg", ".jpeg", ".webp", ".gif"}:
        # Map extension to MIME type for the Anthropic vision API
        _mime_map = {
            ".png": "image/png",
            ".jpg": "image/jpeg",
            ".jpeg": "image/jpeg",
            ".webp": "image/webp",
            ".gif": "image/gif",
        }
        media_type = _mime_map[suffix]
        image_data = base64.standard_b64encode(path.read_bytes()).decode("utf-8")

        try:
            import anthropic as _anthropic
            from archivum.config import get_settings

            settings = get_settings()
            if not settings.anthropic_api_key:
                raise UnsupportedFileTypeError(
                    "ANTHROPIC_API_KEY is not set; cannot parse images without a valid API key"
                )

            client = _anthropic.Anthropic(api_key=settings.anthropic_api_key)
            response = client.messages.create(
                model="claude-haiku-4-5-20251001",
                max_tokens=1024,
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "image",
                                "source": {
                                    "type": "base64",
                                    "media_type": media_type,
                                    "data": image_data,
                                },
                            },
                            {
                                "type": "text",
                                "text": (
                                    "Describe this image in detail. If it contains text, transcribe all text verbatim. "
                                    "Include: what the image shows, any text present, any diagrams or charts described in words."
                                ),
                            },
                        ],
                    }
                ],
            )
            description = response.content[0].text
        except _anthropic.APIError as exc:
            logger.error("Anthropic vision API error for %s: %s", path.name, exc)
            raise UnsupportedFileTypeError(f"Anthropic vision API error: {exc}") from exc

        return ParsedDoc(
            text=description,
            source=str(path),
            metadata={"type": "image", "format": suffix.lstrip(".")},
        )

    # ── Audio ─────────────────────────────────────────────────────────────────
    if suffix in {".mp3", ".m4a", ".wav", ".ogg", ".flac", ".aac", ".aiff", ".opus", ".wma"}:
        model = _get_whisper_model()
        result = model.transcribe(str(path))
        return ParsedDoc(
            text=_format_transcription(result),
            source=str(path),
            metadata=_transcription_metadata(result, "audio"),
        )

    # ── Video ─────────────────────────────────────────────────────────────────
    if suffix in {".mp4", ".mov", ".avi", ".mkv", ".webm", ".m4v", ".mpg", ".mpeg", ".3gp"}:
        model = _get_whisper_model()
        with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as tmp:
            audio_path = tmp.name
        try:
            try:
                subprocess.run(
                    [
                        "ffmpeg", "-i", str(path),
                        "-vn", "-acodec", "pcm_s16le",
                        "-ar", "16000", "-ac", "1",
                        audio_path, "-y",
                    ],
                    check=True,
                    capture_output=True,
                )
            except FileNotFoundError as exc:
                raise UnsupportedFileTypeError(
                    "Video transcription needs ffmpeg. Open Settings and run Install / Enable under Audio Transcription."
                ) from exc
            except subprocess.CalledProcessError as exc:
                stderr = (
                    exc.stderr.decode("utf-8", errors="replace")
                    if isinstance(exc.stderr, bytes)
                    else str(exc.stderr)
                )
                raise UnsupportedFileTypeError(
                    f"ffmpeg could not extract audio from {path.name}: {stderr}"
                ) from exc
            result = model.transcribe(audio_path)
        finally:
            try:
                os.unlink(audio_path)
            except OSError:
                pass
        return ParsedDoc(
            text=_format_transcription(result),
            source=str(path),
            metadata=_transcription_metadata(result, "video"),
        )

    raise UnsupportedFileTypeError(
        f"Archivum cannot parse {suffix or 'this'} files yet. Try a supported document, media, archive, code, data, email, or subtitle file."
    )


async def parse_url(url: str) -> ParsedDoc:
    """Async URL parser — fetches HTML and extracts main content."""
    async with httpx.AsyncClient(
        follow_redirects=True,
        timeout=30.0,
        headers={"User-Agent": "Archivum/1.0 (knowledge-base-bot)"},
    ) as client:
        response = await client.get(url)
        response.raise_for_status()

    content_type = response.headers.get("content-type", "")
    content_type_base = content_type.split(";", 1)[0].strip().lower()

    if "text/html" in content_type or "application/xhtml" in content_type:
        text = _html_to_text(response.text, url=url)
        # Try to extract title
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(response.text, "html.parser")
        title_tag = soup.find("title")
        title = title_tag.text.strip() if title_tag else url

        return ParsedDoc(
            text=text,
            source=url,
            metadata={"type": "url", "url": url, "title": title, "content_type": content_type},
        )

    if "application/json" in content_type:
        return ParsedDoc(
            text=response.text[:8000],
            source=url,
            metadata={"type": "json_url", "url": url},
        )

    if "text/plain" in content_type or "text/markdown" in content_type:
        return ParsedDoc(
            text=response.text,
            source=url,
            metadata={"type": "text_url", "url": url},
        )

    downloadable_types = {
        "application/pdf",
        "application/epub+zip",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/zip",
        "text/rtf",
        "application/rtf",
        "application/xml",
        "text/xml",
    }
    if content_type_base in downloadable_types:
        filename = _download_filename(
            url,
            content_type,
            response.headers.get("content-disposition", ""),
        )
        with tempfile.TemporaryDirectory(prefix="archivum-url-") as tmp:
            tmp_path = Path(tmp) / Path(filename).name
            tmp_path.write_bytes(response.content)
            doc = parse_file(tmp_path)
        return ParsedDoc(
            text=doc.text,
            source=url,
            metadata={
                **doc.metadata,
                "url": url,
                "filename": filename,
                "content_type": content_type,
            },
        )

    # Fallback: try to decode as text
    return ParsedDoc(
        text=response.text[:8000],
        source=url,
        metadata={"type": "unknown_url", "url": url},
    )


async def parse_source(source: str | Path) -> ParsedDoc:
    """Top-level dispatcher: handles Path objects and URL strings."""
    if isinstance(source, str) and (source.startswith("http://") or source.startswith("https://")):
        return await parse_url(source)

    path = Path(source)
    if not path.exists():
        raise FileNotFoundError(f"Source file not found: {path}")

    loop = asyncio.get_running_loop()
    return await loop.run_in_executor(None, parse_file, path)
